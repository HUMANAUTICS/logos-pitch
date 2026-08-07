import os
import subprocess
import sys

def check_and_install_dependencies():
    """Ensure python dependencies are installed."""
    required = {"pymupdf", "python-pptx"}
    installed = set()
    
    # Try importing to see if they exist
    for pkg in required:
        try:
            if pkg == "pymupdf":
                import fitz
            elif pkg == "python-pptx":
                import pptx
            installed.add(pkg)
        except ImportError:
            pass
            
    missing = required - installed
    if missing:
        print(f"Installing missing dependencies: {missing}...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "--break-system-packages", *missing])

def export_to_pdf(html_path, pdf_path):
    """Render HTML to PDF using macOS headless Google Chrome."""
    chrome_path = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
    if not os.path.exists(chrome_path):
        print(f"Error: Google Chrome not found at {chrome_path}", file=sys.stderr)
        sys.exit(1)
        
    print(f"Generating PDF using headless Chrome from {html_path}...")
    abs_html = os.path.abspath(html_path)
    abs_pdf = os.path.abspath(pdf_path)
    
    # Run headless Chrome print-to-pdf command
    cmd = [
        chrome_path,
        "--headless",
        "--disable-gpu",
        "--no-pdf-header-footer",
        f"--print-to-pdf={abs_pdf}",
        abs_html
    ]
    
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"Chrome export failed:\n{result.stderr}", file=sys.stderr)
        sys.exit(1)
        
    print(f"Successfully generated {pdf_path}")

def convert_pdf_to_pptx(pdf_path, pptx_path):
    """Convert PDF pages to a 16:9 PowerPoint presentation as slide images."""
    import fitz  # PyMuPDF
    from pptx import Presentation
    from pptx.util import Inches
    
    print(f"Converting {pdf_path} to PPTX {pptx_path}...")
    
    doc = fitz.open(pdf_path)
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]  # Blank slide layout
    
    for i, page in enumerate(doc):
        print(f"Processing slide {i+1}/{len(doc)}...")
        # Render PDF page to high-quality image (150 DPI)
        pix = page.get_pixmap(dpi=150)
        img_name = f"__temp_slide_{i}.png"
        pix.save(img_name)
        
        # Add slide and insert the image to fill the slide area
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_name, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # Remove temporary image file
        os.remove(img_name)
        
    prs.save(pptx_path)
    print(f"Successfully generated {pptx_path}")

if __name__ == "__main__":
    html_file = "presentation.html"
    pdf_file = "presentation.pdf"
    pptx_file = "presentation.pptx"
    
    check_and_install_dependencies()
    export_to_pdf(html_file, pdf_file)
    convert_pdf_to_pptx(pdf_file, pptx_file)
    print("Export complete!")
